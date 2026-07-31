"""PPTX text extraction for the Content Engine import pipeline."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation as _PptxPresentation


class PptxReader:
    """Extract plain text from PPTX files."""

    def read(self, path: Path) -> str:
        """Read and extract text from all slides of a PPTX file.

        Args:
            path: Path to the PPTX file.

        Returns:
            Extracted text with non-empty shape text joined by newlines, or
            an empty string when no text is available.
        """
        presentation = _PptxPresentation(str(path))
        lines: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", None)
                if text:
                    lines.append(text)
        return "\n".join(lines)
